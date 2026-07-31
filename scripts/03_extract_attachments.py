#!/usr/bin/env python3
"""
Enriquece emails.jsonl añadiendo el texto extraído de cada adjunto
(PDF, DOCX, PPTX, XLSX, TXT, DOC/XLS/PPT legacy, RTF, HTML).

Estrategia de extracción en capas — importante para PDFs, que son el
grueso del corpus: `unstructured` con su estrategia automática puede
devolver 0 elementos EN SILENCIO (sin lanzar excepción) para algunos PDFs
con capa de texto perfectamente válida, aparentemente por cómo resuelve
la detección de layout/GPU en este entorno. Un try/except que solo
reacciona a excepciones NO detecta ese caso. Por eso:

  1. PDF: primero pdftotext (poppler CLI) -> pypdf -> unstructured "fast"
     -> OCR explícito solo si todo lo anterior falla (PDF probablemente
     escaneado). Se descarta también texto "basura" tipo /UNICxxxx que
     algunos PDFs con fuentes subconjunto producen en pypdf.
  2. Resto de formatos: unstructured genérico, con fallback a extractores
     nativos ligeros si unstructured lanza excepción O devuelve vacío.

Uso:
    python 03_extract_attachments.py --emails data/processed/emails.jsonl \
        --output data/processed/emails_enriched.jsonl
"""
import argparse
import json
import re
import subprocess
from pathlib import Path

from tqdm import tqdm

from _jsonl_utils import read_jsonl_lines

SUPPORTED_EXT = {
    ".pdf", ".docx", ".pptx", ".xlsx", ".txt",
    ".doc", ".xls", ".ppt",
    ".rtf", ".htm", ".html",
}

LANGUAGES = ["spa", "eng"]
MIN_MEANINGFUL_CHARS = 20

_UNIC_PLACEHOLDER = re.compile(r"/?UNIC[0-9A-Fa-f]{4}")


def _looks_garbled(text: str) -> bool:
    if not text:
        return True
    matches = _UNIC_PLACEHOLDER.findall(text)
    if not matches:
        return False
    matched_chars = sum(len(m) for m in matches)
    return matched_chars > 0.3 * len(text)


def _is_meaningful(text: str) -> bool:
    return bool(text) and len(text.strip()) >= MIN_MEANINGFUL_CHARS and not _looks_garbled(text)


def extract_via_pdftotext_cli(path: Path) -> str:
    try:
        result = subprocess.run(
            ["pdftotext", "-enc", "UTF-8", str(path), "-"],
            capture_output=True, text=True, timeout=60,
        )
        if result.returncode == 0:
            return result.stdout
    except (FileNotFoundError, subprocess.TimeoutExpired):
        pass
    return ""


def extract_pdf_text(path: Path) -> tuple:
    text = extract_via_pdftotext_cli(path)
    if _is_meaningful(text):
        return text, "pdftotext"

    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
        if _is_meaningful(text):
            return text, "pypdf"
    except Exception:
        pass

    try:
        from unstructured.partition.pdf import partition_pdf
        elements = partition_pdf(filename=str(path), languages=LANGUAGES, strategy="fast")
        text = "\n".join(el.text for el in elements if getattr(el, "text", None))
        if _is_meaningful(text):
            return text, "unstructured_fast"
    except Exception:
        pass

    try:
        from unstructured.partition.pdf import partition_pdf
        elements = partition_pdf(filename=str(path), languages=LANGUAGES, strategy="ocr_only")
        text = "\n".join(el.text for el in elements if getattr(el, "text", None))
        if text:
            return text, "unstructured_ocr"
    except Exception as e:
        return f"[ERROR OCR en {path.name}: {e}]", "error"

    return "", "empty"


def extract_generic_with_unstructured(path: Path) -> str:
    from unstructured.partition.auto import partition
    elements = partition(filename=str(path), languages=LANGUAGES)
    return "\n".join(el.text for el in elements if getattr(el, "text", None))


def extract_fallback(path: Path) -> str:
    ext = path.suffix.lower()
    try:
        if ext == ".docx":
            import docx
            doc = docx.Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs)
        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return "\n".join(texts)
        if ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(str(path), data_only=True)
            texts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    texts.append(" ".join(str(c) for c in row if c is not None))
            return "\n".join(texts)
        if ext == ".txt":
            return path.read_text(errors="replace")
        if ext in (".htm", ".html"):
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(path.read_text(errors="replace"), "html.parser")
            for tag in soup(["script", "style"]):
                tag.decompose()
            return soup.get_text(separator="\n").strip()
    except Exception as e:
        return f"[ERROR extrayendo {path.name}: {e}]"
    return ""


def extract_attachment_text(path_str: str) -> dict:
    path = Path(path_str)
    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXT or not path.exists():
        return {"path": path_str, "text": "", "extractor": "skipped"}

    if ext == ".pdf":
        text, extractor = extract_pdf_text(path)
        return {"path": path_str, "filename": path.name, "text": text, "extractor": extractor}

    try:
        text = extract_generic_with_unstructured(path)
        extractor = "unstructured"
    except Exception:
        text = ""
        extractor = "unstructured_failed"

    if not _is_meaningful(text):
        fallback_text = extract_fallback(path)
        if fallback_text:
            text = fallback_text
            extractor = "fallback"

    return {"path": path_str, "filename": path.name, "text": text, "extractor": extractor}


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--emails", required=True)
    ap.add_argument("--output", required=True)
    args = ap.parse_args()

    in_path = Path(args.emails)
    out_path = Path(args.output)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    lines = read_jsonl_lines(in_path)
    print(f"Procesando {len(lines)} correos...")

    extractor_counts = {}
    with out_path.open("w", encoding="utf-8") as out_f:
        for line in tqdm(lines, desc="Extrayendo adjuntos"):
            record = json.loads(line)
            attachment_texts = []
            for att_path in record.get("attachments", []):
                result = extract_attachment_text(att_path)
                extractor_counts[result["extractor"]] = extractor_counts.get(result["extractor"], 0) + 1
                attachment_texts.append(result)
            record["attachment_texts"] = attachment_texts
            out_f.write(json.dumps(record, ensure_ascii=False) + "\n")

    print(f"Listo. -> {out_path}")
    print("Desglose por método de extracción:")
    for method, count in sorted(extractor_counts.items(), key=lambda x: -x[1]):
        print(f"  {method}: {count}")


if __name__ == "__main__":
    main()
